# -*- coding: utf-8 -*-
"""Generierung von Word-, Excel- und PowerPoint-Dateien aus JSON-Templates."""

from __future__ import annotations

import json
import os
from datetime import datetime
from typing import Any

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt


class OfficeGenerator:
    def __init__(self, templates_path: str, output_dir: str | None = None) -> None:
        self.templates_path = os.path.abspath(templates_path)
        with open(self.templates_path, "r", encoding="utf-8") as f:
            self.templates = json.load(f)
        base = os.path.dirname(self.templates_path)
        self.output_dir = os.path.abspath(output_dir or os.path.join(base, "output"))
        os.makedirs(self.output_dir, exist_ok=True)

    def generate_word_document(
        self,
        template_type: str,
        title: str,
        content: str,
        author: str = "Rambo Rainer",
    ) -> dict[str, Any]:
        word_tpl = self.templates.get("word_templates") or {}
        template = word_tpl.get(template_type)
        if not template:
            return {"error": f"Template {template_type!r} not found"}

        defaults = dict(template.get("default_content") or {})
        doc = Document()

        if defaults.get("header"):
            doc.add_heading(str(defaults["header"]), level=0)

        doc.add_heading(title or "Unbenanntes Dokument", level=1)

        if template_type == "report" and defaults.get("executive_summary"):
            p = doc.add_paragraph()
            r = p.add_run("Zusammenfassung: ")
            r.bold = True
            p.add_run(str(defaults["executive_summary"]))

        doc.add_paragraph(content or "")

        if template_type == "report" and defaults.get("conclusion"):
            doc.add_heading("Fazit", level=2)
            doc.add_paragraph(str(defaults["conclusion"]))

        if defaults.get("signature"):
            doc.add_paragraph(str(defaults["signature"]))

        doc.add_paragraph(
            f"Autor: {author}\nDatum: {datetime.now().strftime('%d.%m.%Y %H:%M')}"
        )

        filename = f"{template_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = os.path.join(self.output_dir, filename)
        doc.save(filepath)

        return {"status": "success", "file": filename, "path": filepath}

    def generate_excel_sheet(
        self,
        template_type: str,
        data: dict[str, Any] | None = None,
        formulas: dict[str, str] | None = None,
    ) -> dict[str, Any]:
        excel_tpl = self.templates.get("excel_templates") or {}
        template = excel_tpl.get(template_type)
        if not template:
            return {"error": f"Template {template_type!r} not found"}

        data = data or {}
        merged_formulas = dict(template.get("formulas") or {})
        if formulas:
            merged_formulas.update(formulas)

        wb = Workbook()
        wb.remove(wb.active)

        for sheet_name in template.get("sheets") or []:
            ws = wb.create_sheet(str(sheet_name))
            ws["A1"] = str(sheet_name)
            ws["A1"].font = Font(bold=True, size=14)
            rows = data.get(sheet_name)
            if isinstance(rows, list):
                for r_idx, row in enumerate(rows, start=2):
                    if isinstance(row, (list, tuple)):
                        for c_idx, val in enumerate(row, start=1):
                            ws.cell(row=r_idx, column=c_idx, value=val)
                    elif isinstance(row, dict):
                        for k, val in row.items():
                            ws[str(k)] = val

        for cell_ref, formula in merged_formulas.items():
            part = str(cell_ref).partition("!")
            if not part[1]:
                continue
            sheet_name, _, cell = part[0], part[1], part[2].strip()
            if sheet_name not in wb.sheetnames:
                continue
            wb[sheet_name][cell] = formula

        filename = f"{template_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(self.output_dir, filename)
        wb.save(filepath)

        return {"status": "success", "file": filename, "path": filepath}

    def generate_powerpoint(
        self,
        template_type: str,
        slides: list[dict[str, Any]] | None = None,
    ) -> dict[str, Any]:
        ppt_tpl = self.templates.get("powerpoint_templates") or {}
        template = ppt_tpl.get(template_type)
        if not template:
            return {"error": f"Template {template_type!r} not found"}

        slide_defs = slides if isinstance(slides, list) and slides else list(template.get("slides") or [])

        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)

        try:
            blank_layout = prs.slide_layouts[6]
        except IndexError:
            blank_layout = prs.slide_layouts[-1]

        for slide_def in slide_defs:
            if not isinstance(slide_def, dict):
                continue
            slide = prs.slides.add_slide(blank_layout)
            stitle = str(slide_def.get("title") or "")

            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1.2)
            )
            tf = title_box.text_frame
            tf.text = stitle
            p0 = tf.paragraphs[0]
            p0.font.size = PptxPt(32)
            p0.font.bold = True

            if slide_def.get("subtitle"):
                p1 = tf.add_paragraph()
                p1.text = str(slide_def["subtitle"])
                p1.font.size = PptxPt(18)

            if slide_def.get("content"):
                content_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(2.0), PptxInches(9), PptxInches(5)
                )
                ctf = content_box.text_frame
                ctf.text = str(slide_def["content"])
                ctf.paragraphs[0].font.size = PptxPt(20)

        filename = f"{template_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)

        return {"status": "success", "file": filename, "path": filepath}
